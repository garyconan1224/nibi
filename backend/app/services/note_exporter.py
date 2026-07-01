"""Single-item NoteShell export helpers."""

from __future__ import annotations

import base64
import io
import mimetypes
import re
import zipfile
from html import escape
from pathlib import Path
from typing import Iterable
from urllib.parse import quote, unquote

from fastapi import HTTPException
from fastapi.responses import StreamingResponse

from backend.app.models.workspace import WorkspaceItem
from shared.config import DATA_DIR


def build_note_export_response(
    *,
    workspace_id: str,
    item_id: str,
    item: WorkspaceItem,
    note_md: str,
    source_md: str = "",
    format: str = "obsidian",
) -> StreamingResponse:
    """Build an export response for one NoteShell note."""

    fmt = format.strip().lower().replace("-", "_")
    body = _strip_frontmatter(note_md)
    title = _title_from_body(body) or item.name or item_id
    safe_title = _safe_filename(title)

    if fmt in {"md", "markdown", "raw_md"}:
        return _stream_bytes(
            note_md.encode("utf-8"),
            "text/markdown; charset=utf-8",
            f"{safe_title}.md",
        )
    if fmt == "obsidian":
        return _build_obsidian_zip(
            title=title,
            safe_title=safe_title,
            workspace_id=workspace_id,
            item_id=item_id,
            item=item,
            body=body,
            source_md=source_md,
        )
    if fmt in {"html", "immersive"}:
        html = _render_note_html(title=title, item=item, body=body)
        return _stream_bytes(
            html.encode("utf-8"),
            "text/html; charset=utf-8",
            f"{safe_title}.html",
        )
    if fmt == "pdf":
        return _build_pdf(title=title, safe_title=safe_title, item=item, body=body)
    if fmt in {"long_image", "png", "image"}:
        return _build_long_image(title=title, safe_title=safe_title, item=item, body=body)
    if fmt in {"docx", "word"}:
        return _build_docx(title=title, safe_title=safe_title, item=item, body=body)
    if fmt in {"pptx", "ppt"}:
        return _build_pptx(title=title, safe_title=safe_title, item=item, body=body)

    raise HTTPException(
        status_code=400,
        detail="unsupported format: "
        f"{format!r}; use md/html/pdf/docx/long_image/pptx/obsidian",
    )


def _strip_frontmatter(markdown: str) -> str:
    if not markdown.startswith("---\n"):
        return markdown
    parts = markdown.split("---\n", 2)
    return parts[2].lstrip() if len(parts) >= 3 else markdown


def _title_from_body(body: str) -> str:
    match = re.search(r"^#\s+(.+)$", body, re.MULTILINE)
    return match.group(1).strip() if match else ""


def _safe_filename(name: str) -> str:
    safe = re.sub(r'[/\\:*?"<>|]+', "_", name).strip("._ ")
    return safe[:96] or "note"


def _stream_bytes(data: bytes, media_type: str, filename: str) -> StreamingResponse:
    return StreamingResponse(
        io.BytesIO(data),
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{quote(filename)}"},
    )


def _build_obsidian_zip(
    *,
    title: str,
    safe_title: str,
    workspace_id: str,
    item_id: str,
    item: WorkspaceItem,
    body: str,
    source_md: str,
) -> StreamingResponse:
    exported_md = f"""---
title: {title}
source: {item.source_value or workspace_id}
workspace_id: {workspace_id}
item_id: {item_id}
tags: [nibi, noteshell]
---

{body}"""

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{safe_title}.md", exported_md)
        if source_md.strip():
            zf.writestr(f"{safe_title}-source.md", source_md)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{quote(safe_title + '-obsidian.zip')}"
        },
    )


def _render_note_html(*, title: str, item: WorkspaceItem, body: str) -> str:
    try:
        import markdown2
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="markdown2 未安装，无法导出 HTML/PDF/长图") from exc

    html_body = markdown2.markdown(
        _embed_markdown_images(body),
        extras=["tables", "fenced-code-blocks", "break-on-newline", "task_list"],
    )
    source_label = _source_label(item.source_value)
    item_type = item.type.upper()
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <title>{escape(title)}</title>
  <style>
    :root {{
      color-scheme: light;
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
      color: #191410;
      background: #f7f3ec;
    }}
    body {{
      margin: 0;
      background: #f7f3ec;
    }}
    .page {{
      width: min(920px, calc(100vw - 72px));
      margin: 0 auto;
      padding: 42px 0 64px;
    }}
    .source {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 16px;
      min-height: 72px;
      margin-bottom: 38px;
      padding: 18px 22px;
      border-radius: 12px;
      color: #f8f5ee;
      background: linear-gradient(135deg, #20242c, #62656c);
    }}
    .source strong {{
      display: block;
      max-width: 680px;
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
      font-size: 16px;
    }}
    .source span {{
      color: rgba(255,255,255,.72);
      font-size: 12px;
    }}
    article {{
      padding: 0 4px;
      font-size: 16px;
      line-height: 1.85;
    }}
    h1 {{
      margin: 0 0 28px;
      padding-bottom: 16px;
      border-bottom: 1px solid rgba(25,20,16,.12);
      font-size: 30px;
      line-height: 1.2;
    }}
    h2 {{
      margin: 34px 0 14px;
      font-size: 23px;
      line-height: 1.25;
    }}
    h3 {{
      margin: 28px 0 10px;
      font-size: 18px;
    }}
    p, li {{
      color: #3d3929;
    }}
    blockquote {{
      margin: 20px 0;
      padding: 14px 18px;
      border-left: 4px solid #e96f22;
      border-radius: 0 8px 8px 0;
      background: #fff0dd;
      color: #2f291d;
    }}
    table {{
      width: 100%;
      border-collapse: collapse;
      margin: 16px 0;
      font-size: 14px;
    }}
    th, td {{
      border: 1px solid rgba(25,20,16,.12);
      padding: 8px 10px;
      vertical-align: top;
    }}
    th {{
      background: rgba(255,255,255,.68);
    }}
    img {{
      max-width: 100%;
      height: auto;
      display: block;
      margin: 18px auto;
      border-radius: 10px;
      box-shadow: 0 12px 30px rgba(28,24,18,.12);
    }}
    code, pre {{
      font-family: "SFMono-Regular", Consolas, monospace;
      background: rgba(25,20,16,.06);
      border-radius: 6px;
    }}
    pre {{
      overflow-x: auto;
      padding: 12px;
    }}
  </style>
</head>
<body>
  <main class="page">
    <section class="source">
      <div>
        <strong>{escape(title)}</strong>
        <span>{escape(item_type)} · {escape(source_label)}</span>
      </div>
      <span>Generated by Nibi</span>
    </section>
    <article>
      <h1>{escape(title)}</h1>
      {html_body}
    </article>
  </main>
</body>
</html>"""


def _source_label(source_value: str) -> str:
    host = (source_value or "").replace("https://", "").replace("http://", "").split("/")[0]
    return host or "本地素材"


def _embed_markdown_images(markdown: str) -> str:
    def replace(match: re.Match[str]) -> str:
        alt = match.group(1)
        src = match.group(2)
        data_uri = _asset_to_data_uri(src)
        return f"![{alt}]({data_uri or src})"

    return re.sub(r"!\[([^\]]*)\]\(([^)]+)\)", replace, markdown)


def _asset_to_data_uri(src: str) -> str:
    path = _resolve_asset_path(src)
    if not path or not path.exists() or not path.is_file():
        return ""
    mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
    return f"data:{mime};base64,{base64.b64encode(path.read_bytes()).decode('ascii')}"


def _resolve_asset_path(src: str) -> Path | None:
    if not src:
        return None
    raw = unquote(src)
    if raw.startswith("/static/"):
        return DATA_DIR / raw.removeprefix("/static/")
    if raw.startswith("static/"):
        return DATA_DIR / raw.removeprefix("static/")
    if raw.startswith(("http://", "https://", "data:")):
        return None
    path = Path(raw)
    if path.is_absolute():
        return path
    return DATA_DIR / path


def _build_pdf(*, title: str, safe_title: str, item: WorkspaceItem, body: str) -> StreamingResponse:
    html = _render_note_html(title=title, item=item, body=body)
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="playwright 未安装，无法导出 PDF") from exc

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1200, "height": 1600})
            page.set_content(html, wait_until="networkidle")
            data = page.pdf(
                format="A4",
                print_background=True,
                margin={"top": "16mm", "right": "16mm", "bottom": "18mm", "left": "16mm"},
            )
            browser.close()
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PDF 导出失败: {exc}") from exc

    return _stream_bytes(data, "application/pdf", f"{safe_title}.pdf")


def _build_long_image(*, title: str, safe_title: str, item: WorkspaceItem, body: str) -> StreamingResponse:
    html = _render_note_html(title=title, item=item, body=body)
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="playwright 未安装，无法导出长图") from exc

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": 1200, "height": 1600}, device_scale_factor=1)
            page.set_content(html, wait_until="networkidle")
            data = page.screenshot(full_page=True, type="png")
            browser.close()
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"长图导出失败: {exc}") from exc

    return _stream_bytes(data, "image/png", f"{safe_title}.png")


def _build_docx(*, title: str, safe_title: str, item: WorkspaceItem, body: str) -> StreamingResponse:
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="python-docx 未安装，无法导出 Word") from exc

    doc = Document()
    doc.add_heading(title, level=0)
    source = doc.add_paragraph()
    source.add_run(f"{item.type.upper()} · {_source_label(item.source_value)}").italic = True

    for line in body.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        img_match = re.match(r"!\[[^\]]*\]\(([^)]+)\)", stripped)
        if img_match:
            path = _resolve_asset_path(img_match.group(1))
            if path and path.exists():
                try:
                    doc.add_picture(str(path), width=Inches(5.6))
                    continue
                except Exception:
                    pass
        if stripped.startswith("# "):
            doc.add_heading(stripped[2:].strip(), level=1)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:].strip(), level=2)
        elif stripped.startswith("### "):
            doc.add_heading(stripped[4:].strip(), level=3)
        elif stripped.startswith(("- ", "* ")):
            doc.add_paragraph(_plain_markdown(stripped[2:]), style="List Bullet")
        elif re.match(r"^\d+\.\s+", stripped):
            doc.add_paragraph(_plain_markdown(re.sub(r"^\d+\.\s+", "", stripped)), style="List Number")
        elif stripped.startswith(">"):
            doc.add_paragraph(_plain_markdown(stripped.lstrip("> ").strip()))
        elif stripped.startswith("|"):
            doc.add_paragraph(_plain_markdown(stripped))
        else:
            doc.add_paragraph(_plain_markdown(stripped))

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return _stream_bytes(
        buf.read(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        f"{safe_title}.docx",
    )


def _build_pptx(*, title: str, safe_title: str, item: WorkspaceItem, body: str) -> StreamingResponse:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="python-pptx 未安装，无法导出 PPT；请先安装 python-pptx",
        ) from exc

    bg = RGBColor(248, 246, 239)
    ink = RGBColor(33, 29, 25)
    muted = RGBColor(111, 104, 92)
    line = RGBColor(222, 214, 199)
    accent = RGBColor(229, 105, 34)
    accent_dark = RGBColor(126, 59, 24)
    slate = RGBColor(32, 36, 44)
    slate_2 = RGBColor(55, 62, 75)
    white = RGBColor(255, 252, 246)
    cream = RGBColor(255, 246, 232)
    sage = RGBColor(220, 232, 215)

    def set_background(slide, color) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = line_color
        return shape

    def add_text(
        slide,
        left,
        top,
        width,
        height,
        text: str,
        *,
        size: int,
        color=ink,
        bold: bool = False,
        align=None,
        vertical=None,
    ):
        shape = slide.shapes.add_textbox(left, top, width, height)
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        if vertical is not None:
            frame.vertical_anchor = vertical
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = "PingFang SC"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        if align is not None:
            paragraph.alignment = align
        return shape

    def add_chip(slide, left, top, text: str, *, fill_color=cream, color=accent_dark, width=None):
        chip_width = width or Inches(max(1.08, min(3.2, 0.34 + len(text) * 0.14)))
        shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, chip_width, Inches(0.34), fill_color, None)
        frame = shape.text_frame
        frame.clear()
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = Inches(0.12)
        frame.margin_right = Inches(0.12)
        p = frame.paragraphs[0]
        p.text = text
        p.font.name = "PingFang SC"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        return shape

    def add_footer(slide, index: int) -> None:
        add_text(slide, Inches(0.62), Inches(7.05), Inches(2.5), Inches(0.22), "Nibi Notes", size=8, color=muted)
        add_text(
            slide,
            Inches(11.55),
            Inches(7.05),
            Inches(1.1),
            Inches(0.22),
            f"{index:02d}",
            size=8,
            color=muted,
            align=PP_ALIGN.RIGHT,
        )

    def add_picture_fit(slide, path: Path, left, top, max_width, max_height) -> bool:
        try:
            frame = add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.05),
                top - Inches(0.05),
                max_width + Inches(0.1),
                max_height + Inches(0.1),
                white,
                line,
            )
            frame.shadow.inherit = False
            picture = slide.shapes.add_picture(str(path), left, top, width=max_width)
            if picture.height > max_height:
                scale = max_height / picture.height
                picture.width = int(picture.width * scale)
                picture.height = int(max_height)
            if picture.width > max_width:
                scale = max_width / picture.width
                picture.width = int(max_width)
                picture.height = int(picture.height * scale)
            picture.left = int(left + (max_width - picture.width) / 2)
            picture.top = int(top + (max_height - picture.height) / 2)
            return True
        except Exception:
            return False

    def add_bullets(slide, left, top, width, height, points: list[str], *, size: int = 17, color=ink) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        for idx, point in enumerate(points):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = point
            paragraph.level = 0
            paragraph.font.name = "PingFang SC"
            paragraph.font.size = Pt(size)
            paragraph.font.color.rgb = color
            paragraph.space_after = Pt(10)

    def add_table(slide, rows: list[list[str]]) -> None:
        if len(rows) < 2:
            return
        col_count = min(4, max(len(row) for row in rows))
        normalized = [(row + [""] * col_count)[:col_count] for row in rows[:7]]
        table_shape = slide.shapes.add_table(
            len(normalized),
            col_count,
            Inches(0.9),
            Inches(1.85),
            Inches(11.55),
            Inches(4.55),
        )
        table = table_shape.table
        for col_idx in range(col_count):
            table.columns[col_idx].width = int(Inches(11.55) / col_count)
        for row_idx, row in enumerate(normalized):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = _ppt_truncate(value, 54)
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
                cell.fill.solid()
                cell.fill.fore_color.rgb = slate if row_idx == 0 else white
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = "PingFang SC"
                    paragraph.font.size = Pt(11 if row_idx == 0 else 10)
                    paragraph.font.bold = row_idx == 0
                    paragraph.font.color.rgb = white if row_idx == 0 else ink

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    sections = _ppt_sections(body)[:10]
    if not sections:
        sections = [{
            "title": "内容概览",
            "timestamp": "",
            "points": [_ppt_truncate(_plain_markdown(body), 150) or "本节主要内容见原笔记"],
            "image_path": _first_image_path(body.splitlines()),
            "table": [],
        }]

    cover_image = next((section["image_path"] for section in sections if section["image_path"]), None)
    source_line = f"{item.type.upper()} · {_source_label(item.source_value)}"

    slide_no = 1
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(cover, slate)
    add_shape(cover, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), slate, None)
    add_shape(cover, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5), accent, None)
    add_chip(cover, Inches(0.75), Inches(0.72), "NOTE DECK", fill_color=RGBColor(255, 236, 216), color=accent_dark)
    add_text(cover, Inches(0.75), Inches(1.35), Inches(6.8), Inches(2.0), _ppt_truncate(title, 58), size=34, color=white, bold=True)
    add_text(cover, Inches(0.78), Inches(3.62), Inches(5.8), Inches(0.38), source_line, size=12, color=RGBColor(209, 214, 222))
    add_text(
        cover,
        Inches(0.78),
        Inches(4.18),
        Inches(5.65),
        Inches(1.25),
        "从原始笔记结构自动整理：封面、目录、章节要点、关键图片与表格。",
        size=17,
        color=RGBColor(235, 229, 219),
    )
    if cover_image:
        add_picture_fit(cover, cover_image, Inches(8.05), Inches(1.08), Inches(4.45), Inches(4.95))
    else:
        add_shape(cover, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.18), Inches(4.2), Inches(4.55), slate_2, None)
        add_text(cover, Inches(8.65), Inches(2.58), Inches(3.3), Inches(0.95), "Structured\nNotes", size=27, color=white, bold=True, align=PP_ALIGN.CENTER)
    add_footer(cover, slide_no)

    slide_no += 1
    agenda = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(agenda, bg)
    add_chip(agenda, Inches(0.72), Inches(0.62), "目录", fill_color=cream, color=accent_dark)
    add_text(agenda, Inches(0.72), Inches(1.08), Inches(5.8), Inches(0.55), "本次笔记结构", size=27, color=ink, bold=True)
    add_shape(agenda, MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.78), Inches(11.9), Inches(0.02), accent, None)
    for idx, section in enumerate(sections[:8], start=1):
        row_top = Inches(1.98 + (idx - 1) * 0.56)
        add_text(agenda, Inches(0.88), row_top, Inches(0.48), Inches(0.28), f"{idx:02d}", size=11, color=accent_dark, bold=True)
        heading = section["title"]
        if section["timestamp"]:
            heading = f"{heading}  [{section['timestamp']}]"
        add_text(agenda, Inches(1.55), row_top - Inches(0.03), Inches(8.6), Inches(0.34), _ppt_truncate(heading, 52), size=16, color=ink)
    if len(sections) > 8:
        add_text(agenda, Inches(1.55), Inches(6.56), Inches(8.6), Inches(0.32), f"另外 {len(sections) - 8} 个章节详见原笔记", size=13, color=muted)
    add_footer(agenda, slide_no)

    for section in sections:
        slide_no += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, bg)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18), accent, None)
        if section["timestamp"]:
            add_chip(slide, Inches(0.72), Inches(0.62), f"原片 {section['timestamp']}", fill_color=sage, color=RGBColor(52, 92, 70))
        add_text(slide, Inches(0.72), Inches(1.04), Inches(7.1), Inches(0.9), _ppt_truncate(section["title"], 42), size=25, color=ink, bold=True)
        points = section["points"] or ["本节主要内容见原笔记"]
        has_image = bool(section["image_path"])
        add_bullets(slide, Inches(0.86), Inches(2.22), Inches(6.8), Inches(3.8), points[:5], size=17)
        if has_image:
            add_picture_fit(slide, section["image_path"], Inches(8.1), Inches(1.58), Inches(4.25), Inches(4.65))
        else:
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(2.15), Inches(3.3), Inches(2.85), cream, None)
            add_text(
                slide,
                Inches(8.85),
                Inches(3.08),
                Inches(2.5),
                Inches(0.88),
                "Key\nTakeaway",
                size=23,
                color=accent_dark,
                bold=True,
                align=PP_ALIGN.CENTER,
                vertical=MSO_ANCHOR.MIDDLE,
            )
        add_footer(slide, slide_no)

        if section["table"]:
            slide_no += 1
            table_slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_background(table_slide, bg)
            add_chip(table_slide, Inches(0.72), Inches(0.62), "结构化表格", fill_color=cream, color=accent_dark)
            add_text(table_slide, Inches(0.72), Inches(1.05), Inches(9.3), Inches(0.55), _ppt_truncate(section["title"], 46), size=24, color=ink, bold=True)
            add_table(table_slide, section["table"])
            add_footer(table_slide, slide_no)

    slide_no += 1
    close = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(close, slate)
    add_shape(close, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), slate, None)
    add_chip(close, Inches(0.78), Inches(0.75), "TAKEAWAYS", fill_color=RGBColor(255, 236, 216), color=accent_dark)
    add_text(close, Inches(0.78), Inches(1.35), Inches(5.6), Inches(0.7), "可行动结论", size=30, color=white, bold=True)
    takeaway_points = _ppt_takeaways(sections)
    add_bullets(close, Inches(0.9), Inches(2.35), Inches(10.9), Inches(2.9), takeaway_points, size=18, color=white)
    add_text(close, Inches(0.9), Inches(6.2), Inches(4.8), Inches(0.32), source_line, size=10, color=RGBColor(209, 214, 222))
    add_footer(close, slide_no)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return _stream_bytes(
        buf.read(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"{safe_title}.pptx",
    )


def _ppt_sections(body: str) -> list[dict[str, object]]:
    sections: list[dict[str, object]] = []
    for raw_title, lines in _split_markdown_sections(body):
        title, timestamp = _ppt_clean_title(raw_title)
        points = _ppt_bullet_lines(lines)
        table = _first_markdown_table(lines)
        image_path = _first_image_path(lines)
        if not points and not table and not image_path:
            continue
        sections.append({
            "title": title,
            "timestamp": timestamp,
            "points": points,
            "image_path": image_path,
            "table": table,
        })
    return sections


def _ppt_clean_title(title: str) -> tuple[str, str]:
    timestamp_match = re.search(r"(?:原片\s*@\s*)?(\d{1,2}:\d{2}(?::\d{2})?)", title)
    timestamp = timestamp_match.group(1) if timestamp_match else ""
    cleaned = re.sub(r"\s*(?:原片\s*@\s*)?\d{1,2}:\d{2}(?::\d{2})?\s*", " ", title)
    cleaned = cleaned.strip(" -—[]")
    return cleaned or title, timestamp


def _ppt_bullet_lines(lines: list[str]) -> list[str]:
    points: list[str] = []
    seen: set[str] = set()
    in_code = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code or not stripped:
            continue
        if stripped.startswith(("!", "|", "#")):
            continue
        if re.match(r"^\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?$", stripped):
            continue
        cleaned = re.sub(r"^\s*(?:[-*+]|\d+[.)]|[•◦])\s+", "", stripped)
        cleaned = _plain_markdown(cleaned)
        if not cleaned:
            continue
        cleaned = _ppt_truncate(cleaned, 108)
        key = cleaned.lower()
        if key in seen:
            continue
        points.append(cleaned)
        seen.add(key)
        if len(points) >= 5:
            break
    return points


def _first_markdown_table(lines: list[str]) -> list[list[str]]:
    for index, line in enumerate(lines[:-1]):
        if not line.strip().startswith("|"):
            continue
        separator = lines[index + 1].strip()
        if not re.match(r"^\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?$", separator):
            continue
        rows: list[list[str]] = []
        cursor = index
        while cursor < len(lines) and lines[cursor].strip().startswith("|"):
            if cursor == index + 1:
                cursor += 1
                continue
            cells = [
                _plain_markdown(cell).strip()
                for cell in lines[cursor].strip().strip("|").split("|")
            ]
            if any(cells):
                rows.append(cells)
            cursor += 1
        if len(rows) >= 2:
            return rows
    return []


def _ppt_takeaways(sections: list[dict[str, object]]) -> list[str]:
    points: list[str] = []
    for section in sections:
        section_points = section.get("points")
        if not isinstance(section_points, list):
            continue
        for point in section_points:
            if isinstance(point, str) and point:
                points.append(point)
            if len(points) >= 4:
                return points
    return points or ["保留原笔记结构，继续补充下一步行动。"]


def _ppt_truncate(text: str, limit: int) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 1)].rstrip() + "…"


def _plain_markdown(text: str) -> str:
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"[*_`>#]", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _split_markdown_sections(markdown: str) -> Iterable[tuple[str, list[str]]]:
    current_title = "内容概览"
    current_lines: list[str] = []
    for line in markdown.splitlines():
        match = re.match(r"^(#{2,3})\s+(.+)", line)
        if match:
            if current_lines:
                yield current_title, current_lines
            current_title = _plain_markdown(match.group(2)) or "未命名章节"
            current_lines = []
            continue
        if line.startswith("# "):
            continue
        current_lines.append(line)
    if current_lines:
        yield current_title, current_lines


def _first_image_path(lines: list[str]) -> Path | None:
    for line in lines:
        match = re.search(r"!\[[^\]]*\]\(([^)]+)\)", line)
        if match:
            path = _resolve_asset_path(match.group(1))
            if path and path.exists():
                return path
    return None
